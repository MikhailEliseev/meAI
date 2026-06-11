"""Telegram Gateway — Bot API webhook (incoming) + getUpdates polling (fallback) + Telethon (outgoing via tools).

Per D-16: Hybrid architecture:
- getUpdates polling: incoming messages FROM clients (primary — server can't receive inbound 443 from Telegram)
- Bot API webhook: fallback when direct connectivity works
- Telethon user-client: outgoing messages AS Mikhail, channel search, monitoring

Per D-17: Unified chat — one Operator serves both web chat and Telegram.
Messages from Telegram flow through the same Hermes AIAgent as web chat.

Per D-18: Session binding via tg:// deep link.
Client clicks button in web chat → opens Telegram with deep link containing
web_session_id → bot receives /start command with web_session_id parameter
→ bot links Telegram chat_id to AIM lead dossier.
"""

import asyncio
import json
import logging
import os
import threading
import time
from concurrent.futures import Future
from pathlib import Path
from typing import Optional

import httpx
from fastapi import APIRouter, HTTPException, Request
from pydantic import BaseModel

from .voice_transcriber import transcribe_voice

logger = logging.getLogger(__name__)

TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN", "")
_TELEGRAM_ADMIN_CHAT_IDS: set[int] = set(
    int(x.strip()) for x in os.getenv("TELEGRAM_ADMIN_CHAT_ID", "").split(",")
    if x.strip().isdigit()
)


def _get_mode(chat_id: int, lead_id: str | None) -> str:
    """Determine client mode for Telegram messages.

    Admin chats always get ADMIN mode with full tool access.
    Active leads get ACTIVE, new chats get PRESALE.
    """
    if chat_id in _TELEGRAM_ADMIN_CHAT_IDS:
        return "ADMIN"
    if lead_id:
        return "ACTIVE"
    return "PRESALE"

# Telegram API proxy — hosting in NL blocks Telegram IPs on port 443.
# Old OmniRoute server at 193.111.152.14 acts as HTTP forward proxy.
TELEGRAM_PROXY_URL = os.getenv("TELEGRAM_PROXY_URL", "http://193.111.152.14:7451")
TELEGRAM_PROXY_AUTH = os.getenv("TELEGRAM_PROXY_AUTH", "U9pjtK:hxtlqz")


def _get_proxy_url() -> str | None:
    """Build Telegram proxy URL. Returns None if not configured."""
    if TELEGRAM_PROXY_URL and TELEGRAM_PROXY_AUTH:
        from urllib.parse import urlparse
        parsed = urlparse(TELEGRAM_PROXY_URL)
        host = parsed.netloc or parsed.path
        return f"http://{TELEGRAM_PROXY_AUTH}@{host}"
    return None


router = APIRouter(prefix="/telegram", tags=["telegram"])

# In-memory session binding store (move to DB later)
# Maps web_session_id -> AIM lead_id
_session_bindings: dict[str, str] = {}

# Maps telegram chat_id -> AIM lead_id
_chat_lead_map: dict[int, str] = {}

# Persistent binding storage
_BINDINGS_FILE = Path("/opt/data/chat_bindings.json")


def _load_bindings() -> dict[int, str]:
    """Load chat→lead bindings from disk. Called on module import."""
    try:
        if _BINDINGS_FILE.exists():
            data = json.loads(_BINDINGS_FILE.read_text())
            # JSON keys are strings, convert to int
            bindings = {int(k): v for k, v in data.items()}
            logger.info(f"Loaded {len(bindings)} chat bindings from {_BINDINGS_FILE}")
            return bindings
    except Exception as e:
        logger.warning(f"Failed to load bindings from {_BINDINGS_FILE}: {e}")
    return {}


def _save_bindings() -> None:
    """Persist chat→lead bindings to disk. Called on every change."""
    try:
        _BINDINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
        # JSON keys must be strings
        data = {str(k): v for k, v in _chat_lead_map.items()}
        _BINDINGS_FILE.write_text(json.dumps(data, ensure_ascii=False, indent=2))
    except Exception as e:
        logger.warning(f"Failed to save bindings to {_BINDINGS_FILE}: {e}")


# Load persisted bindings at startup
_chat_lead_map = _load_bindings()

# Polling control
_polling_future: Optional[Future] = None
_polling_stop = False
_last_update_id: int = 0


# ── Models ──────────────────────────────────────────────────────────
class TelegramUpdate(BaseModel):
    update_id: int
    message: Optional[dict] = None
    callback_query: Optional[dict] = None


class DeepLinkBindRequest(BaseModel):
    web_session_id: str
    lead_id: str


class ChatBindRequest(BaseModel):
    chat_id: int
    lead_id: str


# ── Deep Link Binding ───────────────────────────────────────────────
# D-18: tg:// deep link from website binds Telegram chat to AIM lead

@router.post("/bind-session")
async def bind_session(body: DeepLinkBindRequest):
    """Called by Next.js when user clicks 'Open in Telegram' button.

    Creates a pending binding: when the Telegram bot receives a /start
    command with matching web_session_id, it links the Telegram chat_id
    to the AIM lead_id.
    """
    _session_bindings[body.web_session_id] = body.lead_id
    logger.info(f"Session binding created: {body.web_session_id} -> {body.lead_id}")
    return {"status": "bound", "web_session_id": body.web_session_id}


@router.post("/bind-chat")
async def bind_chat(body: ChatBindRequest):
    """Manually bind a Telegram chat_id to an AIM lead_id.

    Admin endpoint — allows direct chat→lead binding without deep link flow.
    Persists to /opt/data/chat_bindings.json to survive restarts.
    """
    _chat_lead_map[body.chat_id] = body.lead_id
    _save_bindings()
    logger.info(f"Manual chat binding: chat_id={body.chat_id} -> lead_id={body.lead_id}")
    return {"status": "bound", "chat_id": body.chat_id, "lead_id": body.lead_id}


@router.get("/list-chats")
async def list_chats():
    """List all bound Telegram chats. Admin debug endpoint."""
    return {
        "bindings": {str(k): v for k, v in _chat_lead_map.items()},
        "pending_sessions": dict(_session_bindings),
        "total": len(_chat_lead_map),
    }


# ── Webhook ─────────────────────────────────────────────────────────
# D-16: Bot API webhook receives incoming messages from clients

@router.post("/webhook")
async def telegram_webhook(request: Request):
    """Telegram Bot API webhook endpoint.

    Receives incoming messages from clients via Telegram Bot API.
    Routes them through the same Hermes AIAgent as web chat (D-17).

    Webhook URL: https://iamaim.ru/telegram/webhook
    Set via: curl -X POST https://api.telegram.org/bot<TOKEN>/setWebhook
             -d url=https://iamaim.ru/telegram/webhook
    """
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON")

    update = TelegramUpdate(**body)

    # Handle /start command with deep link (D-18)
    if update.message:
        text = update.message.get("text", "")
        chat_id = update.message.get("chat", {}).get("id")

        if text.startswith("/start ") and chat_id:
            # Extract web_session_id from deep link
            # Format: /start bind_<web_session_id>
            parts = text.split()
            if len(parts) > 1:
                param = parts[1]
                if param.startswith("bind_"):
                    web_session_id = param[5:]  # Remove "bind_" prefix
                    lead_id = _session_bindings.pop(web_session_id, None)
                    if lead_id:
                        _chat_lead_map[chat_id] = lead_id
                        _save_bindings()
                        logger.info(f"Session bound: chat_id={chat_id} -> lead_id={lead_id}")
                        await _send_telegram_message(
                            chat_id,
                            "✅ Ваш чат привязан к аккаунту AIM. Оператор готов ответить на ваши вопросы."
                        )
                        return {"status": "bound", "lead_id": lead_id}

        # Detect voice message — download + transcribe before processing
        voice = update.message.get("voice")
        if chat_id and voice and not text:
            file_id = voice.get("file_id")
            if file_id:
                await _send_telegram_message(chat_id, "🎤 Расшифровываю голосовое сообщение...")
                text = transcribe_voice(file_id)
                if not text:
                    await _send_telegram_message(chat_id, "❌ Не удалось расшифровать голосовое сообщение")
                    return {"status": "transcription_failed"}

        # Detect document/file — download + extract text (PPTX, PDF, DOCX, TXT)
        document = update.message.get("document")
        if chat_id and document and not text:
            file_id = document.get("file_id")
            file_name = document.get("file_name", "file.bin")
            caption = update.message.get("caption", "")
            if file_id:
                await _send_telegram_message(chat_id, f"📎 Читаю файл «{file_name}»...")
                extracted = await _download_and_extract_text(file_id, file_name)
                if extracted:
                    text = f"[Файл: {file_name}]\n\n{extracted}"
                    if caption:
                        text = f"[Задача от пользователя]: {caption}\n\n{text}"
                elif caption:
                    text = f"[Пользователь отправил файл «{file_name}» с подписью]: {caption}"
                    await _send_telegram_message(
                        chat_id,
                        f"⚠️ Не удалось извлечь текст из «{file_name}». Использую только подпись к файлу."
                    )

        # Detect photo — use caption as message
        photo = update.message.get("photo")
        if chat_id and photo and not text:
            caption = update.message.get("caption", "")
            if caption:
                text = f"[Пользователь отправил фото с подписью]: {caption}"
            else:
                await _send_telegram_message(chat_id, "📸 Я вижу фото. Напишите, что с ним нужно сделать?")
                return {"status": "photo_no_caption"}

        # Process message via Hermes AIAgent with session memory (D-17: unified chat)
        # ASYNC: webhook returns 200 immediately, agent runs in background.
        # Long tool calls (prescan up to 300s) would cause Telegram webhook
        # timeout (~30s) and retry loops if processed synchronously.
        if chat_id and text:
            # Block system-generated messages to prevent self-reply loops.
            # These patterns leak from error handlers and webhook retries.
            if _is_system_message(text):
                logger.info(f"Ignoring system message in chat {chat_id}: {text[:100]}")
                return {"status": "ignored_system_message"}
            asyncio.create_task(_process_message_async(chat_id, text))
            return {"status": "processing"}

    return {"status": "ok"}


# ── Helpers ─────────────────────────────────────────────────────────

_SYSTEM_MESSAGE_PATTERNS = [
    "⚠️ The model provider failed after retries",
    "## SELF-LEARNING — выполни ОБЯЗАТЕЛЬНО",
    "⚠️ Iteration budget exhausted",
    "⚠️ Произошла ошибка при обработке",
    "The model provider failed",
    "⚠️ Reached maximum iterations",
]


def _is_system_message(text: str) -> bool:
    """Detect system-generated messages that would create self-reply loops."""
    for pattern in _SYSTEM_MESSAGE_PATTERNS:
        if pattern in text:
            return True
    return False


async def _download_and_extract_text(file_id: str, file_name: str) -> str | None:
    """Download a file from Telegram Bot API and extract text content.

    Supports: .pptx, .pdf, .docx, .txt, .md, .py, .json, .html, .csv
    Returns extracted text or None if extraction fails.
    """
    import io
    from pathlib import Path

    ext = Path(file_name).suffix.lower()
    get_url = f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/getFile?file_id={file_id}"
    proxy = os.getenv("TELEGRAM_PROXY_URL", "")
    proxy_auth = os.getenv("TELEGRAM_PROXY_AUTH", "")
    proxy_url = None
    if proxy and proxy_auth:
        from urllib.parse import urlparse
        p = urlparse(proxy)
        proxy_url = f"http://{proxy_auth}@{p.netloc or p.path}"

    try:
        async with httpx.AsyncClient(timeout=30.0, proxy=proxy_url) as client:
            resp = await client.get(get_url)
            if resp.status_code != 200:
                logger.error(f"getFile failed: HTTP {resp.status_code}")
                return None
            file_info = resp.json()
            if not file_info.get("ok"):
                logger.error(f"getFile API error: {file_info.get('description')}")
                return None

            file_path = file_info["result"]["file_path"]
            download_url = f"https://api.telegram.org/file/bot{TELEGRAM_BOT_TOKEN}/{file_path}"
            resp = await client.get(download_url)
            if resp.status_code != 200:
                logger.error(f"File download failed: HTTP {resp.status_code}")
                return None

            content = resp.content
            logger.info(f"Downloaded {file_name} ({len(content)} bytes, ext={ext})")

    except Exception as e:
        logger.error(f"Failed to download file {file_name}: {e}")
        return None

    # Extract text based on file type
    try:
        if ext in (".txt", ".md", ".py", ".json", ".html", ".csv", ".yaml", ".yml", ".xml"):
            return content.decode("utf-8", errors="replace")

        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                parts.append(t)
                if parts:
                    slides.append(f"Слайд {i}:\n" + "\n".join(parts))
            return "\n\n".join(slides) if slides else None

        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(content))
            pages = []
            for i, page in enumerate(reader.pages, 1):
                t = page.extract_text()
                if t and t.strip():
                    pages.append(f"Страница {i}:\n{t.strip()}")
            return "\n\n".join(pages) if pages else None

        if ext == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(content))
            paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paras) if paras else None

        logger.warning(f"No text extractor for extension {ext}")
        return None

    except Exception as e:
        logger.error(f"Text extraction failed for {file_name} ({ext}): {e}")
        return None


async def _process_message_async(chat_id: int, text: str) -> None:
    """Process a Telegram message in background and send reply via sendMessage.

    Offloads the synchronous _call_hermes_agent to a thread pool so the
    webhook can return 200 immediately. Long tool calls (prescan,
    competitors) would otherwise hit Telegram's ~30s webhook timeout.
    """
    lead_id = _chat_lead_map.get(chat_id)
    mode = _get_mode(chat_id, lead_id)

    logger.info(f"Processing tg message (async): chat_id={chat_id} mode={mode} text={text[:80]}")
    stop_typing = threading.Event()
    typing_thread = threading.Thread(
        target=_send_chat_action_keepalive,
        args=(chat_id, stop_typing),
        daemon=True,
    )
    typing_thread.start()
    try:
        loop = asyncio.get_running_loop()
        reply = await loop.run_in_executor(
            None, _call_hermes_agent, mode, text, f"tg:{chat_id}"
        )
        logger.info(f"Agent reply received (async): {len(str(reply))} chars. Sending to chat {chat_id}...")
        await _send_telegram_message(chat_id, str(reply))
    except Exception as e:
        logger.error(f"Agent processing failed for chat {chat_id}: {e}", exc_info=True)
        await _send_telegram_message(
            chat_id,
            "⚠️ Произошла ошибка при обработке сообщения. Пожалуйста, попробуйте ещё раз или свяжитесь с администратором."
        )
    finally:
        stop_typing.set()


async def _send_telegram_message(chat_id: int, text: str) -> dict:
    """Send message via Bot API. Async version for webhook."""
    return _send_telegram_message_sync(chat_id, text)


def _send_chat_action_sync(chat_id: int, action: str = "typing") -> bool:
    """Send chat action to Telegram — shows 'typing...' indicator.

    Returns True on success. Telegram auto-expires the indicator after ~5s.
    Use _send_chat_action_keepalive() to keep it alive during long operations.
    """
    if not TELEGRAM_BOT_TOKEN:
        return False
    url = f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/sendChatAction"
    proxy = _get_proxy_url()
    try:
        with httpx.Client(timeout=5.0, proxy=proxy) as client:
            resp = client.post(url, json={"chat_id": chat_id, "action": action})
            data = resp.json()
            if not data.get("ok"):
                logger.warning(f"sendChatAction failed: {data}")
                return False
            return True
    except Exception as e:
        logger.warning(f"sendChatAction error: {e}")
        return False


def _send_chat_action_keepalive(chat_id: int, stop_event: threading.Event, interval: float = 4.0):
    """Keep typing indicator alive by re-sending every `interval` seconds.

    Runs in a daemon thread. Telegram expires typing after ~5s,
    so we refresh at 4s intervals. Stops when stop_event is set.
    """
    while not stop_event.wait(interval):
        _send_chat_action_sync(chat_id)


def _send_telegram_message_sync(chat_id: int, text: str, retries: int = 3) -> dict:
    """Send message via Bot API — synchronous, with retry + HTML→plaintext fallback.

    Tries HTML first (rich formatting). If Telegram rejects the HTML
    (e.g. malformed tags from _apply_markdown_formatting), falls back
    to plain text automatically.
    """
    if not TELEGRAM_BOT_TOKEN:
        logger.error("TELEGRAM_BOT_TOKEN not configured")
        return {"error": "Bot token not configured"}

    url = f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/sendMessage"
    proxy = _get_proxy_url()

    # If report is too large for a single Telegram message, write to file first
    truncated = text[:4096]
    if len(text) > 3500:
        import datetime as _dt
        ts = _dt.datetime.now().strftime("%Y-%m-%d-%H%M%S")
        filepath = f"/opt/data/reports/report-{ts}.md"
        try:
            import os as _os
            _os.makedirs("/opt/data/reports", exist_ok=True)
            with open(filepath, "w") as f:
                f.write(text)
            logger.info(f"Full report saved to {filepath} ({len(text)} chars)")
            truncated = (
                f"📄 Полный отчёт ({len(text)} символов) сохранён: `{filepath}`\n\n"
                + truncated[:3500]
            )
        except Exception as e:
            logger.warning(f"Failed to save report to file: {e}")

    # Try HTML first, fall back to plain text on parse errors
    for parse_mode in ("HTML", None):
        last_error = None
        for attempt in range(retries):
            try:
                body = {"chat_id": chat_id, "text": truncated}
                if parse_mode:
                    body["parse_mode"] = parse_mode
                with httpx.Client(timeout=10.0, proxy=proxy) as client:
                    response = client.post(url, json=body)
                    result = response.json()
                    if result.get("ok"):
                        return result
                    last_error = result
                    err_desc = str(result.get("description", ""))[:120]
                    # HTML parse error — don't retry with same mode, jump to plaintext
                    if "can't parse entities" in err_desc:
                        logger.warning(
                            f"sendMessage HTML parse error, falling back to plaintext: {err_desc}"
                        )
                        break
                    logger.warning(
                        f"sendMessage Telegram error (attempt {attempt+1}/{retries}): {err_desc}"
                    )
                    if attempt < retries - 1:
                        time.sleep(1 * (attempt + 1))
                    continue
            except Exception as e:
                logger.error(f"sendMessage failed (attempt {attempt+1}/{retries}): {e}")
                last_error = str(e)
                if attempt < retries - 1:
                    time.sleep(1 * (attempt + 1))
        else:
            # All retries exhausted for this parse_mode
            if parse_mode == "HTML":
                logger.warning("HTML mode failed, trying plaintext...")
                continue  # try next parse_mode
            logger.error(f"sendMessage FAILED after {retries} retries: {last_error}")
            return {"error": str(last_error)}
        # If we broke out of retry loop due to HTML parse error, continue to plaintext
        continue

    logger.error(f"sendMessage FAILED: both HTML and plaintext failed")
    return {"error": "all modes failed"}


async def _get_updates(offset: int = 0, timeout: int = 30) -> list[dict]:
    """Fetch pending updates from Telegram via long-polling. Async version for webhook."""
    return _get_updates_sync(offset, timeout)


def _get_updates_sync(offset: int = 0, timeout: int = 30) -> list[dict]:
    """Fetch pending updates from Telegram via long-polling — synchronous, for thread."""
    if not TELEGRAM_BOT_TOKEN:
        return []
    url = f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/getUpdates"
    proxy = _get_proxy_url()
    try:
        with httpx.Client(timeout=float(timeout + 10), proxy=proxy) as client:
            response = client.post(url, json={
                "offset": offset,
                "timeout": timeout,
                "allowed_updates": ["message", "callback_query"],
            })
            data = response.json()
            if data.get("ok"):
                return data.get("result", [])
            logger.error(f"getUpdates error: {data}")
            return []
    except Exception as e:
        logger.error(f"getUpdates failed: {e}")
        return []


def _process_update_sync(message_data: dict, chat_id: int, text: str):
    """Process update via Hermes AIAgent — session, tools, memory. Polling thread."""
    lead_id = _chat_lead_map.get(chat_id)
    mode = _get_mode(chat_id, lead_id)

    logger.info(f"Processing tg message: chat_id={chat_id} mode={mode} text={text[:80]}")
    stop_typing = threading.Event()
    typing_thread = threading.Thread(
        target=_send_chat_action_keepalive,
        args=(chat_id, stop_typing),
        daemon=True,
    )
    typing_thread.start()
    try:
        reply = _call_hermes_agent(mode=mode, user_message=text, session_id=f"tg:{chat_id}")
    finally:
        stop_typing.set()
    logger.info(f"Agent reply received: {len(reply)} chars. Sending to Telegram chat {chat_id}...")
    try:
        result = _send_telegram_message_sync(chat_id, reply)
        ok = result.get('ok', False)
        msg_id = 'N/A'
        if isinstance(result.get('result'), dict):
            msg_id = result['result'].get('message_id', 'N/A')
        logger.info(f"sendMessage result: ok={ok} msg_id={msg_id}")
    except Exception as e:
        logger.error(f"CRITICAL: Failed to process or send reply: {e}", exc_info=True)
        try:
            _send_telegram_message_sync(chat_id, f"❌ Ошибка при обработке: {str(e)[:200]}")
        except Exception:
            pass


def _call_hermes_agent(mode: str, user_message: str, session_id: str) -> str:
    """Process message via Hermes AIAgent — full session, tools, SOUL.md identity.

    Uses AIAgent (not raw OmniRoute) so each Telegram chat gets:
    - Session memory (SQLite persistence per chat_id)
    - Tool access (run_seo_audit, collect_contact, etc.)
    - Full SOUL.md identity via load_soul_identity
    - Consistent behavior with web chat
    """
    from .agent_wrapper import run_agent_sync

    result = run_agent_sync(
        message=user_message,
        session_id=session_id,
        mode=mode,
    )
    reply = result.get("reply", "")
    if isinstance(reply, dict):
        reply = reply.get("response", reply.get("content", str(reply)))
    return str(reply)


def _polling_loop_sync():
    """Background thread: poll Telegram for updates via getUpdates.

    Runs in a separate OS thread via run_in_executor to avoid
    httpx proxy connections blocking the uvicorn event loop.

    Uses long-polling (timeout=30s) to reduce request count.
    Reconnects on error with 5s backoff.
    """
    global _last_update_id, _polling_stop

    logger.info("Telegram polling started (getUpdates, sync thread)")
    consecutive_errors = 0
    poll_count = 0

    while not _polling_stop:
        try:
            poll_count += 1
            if poll_count == 1 or poll_count % 20 == 0:
                logger.info(f"Polling alive: poll_count={poll_count} last_update_id={_last_update_id}")
            updates = _get_updates_sync(offset=_last_update_id + 1, timeout=30)
            consecutive_errors = 0

            for update in updates:
                update_id = update.get("update_id", 0)
                if update_id > _last_update_id:
                    _last_update_id = update_id

                message = update.get("message")
                if not message:
                    continue

                text = message.get("text", "")
                chat = message.get("chat", {})
                chat_id = chat.get("id")

                # Detect voice message — download + transcribe before processing
                voice = message.get("voice")
                if chat_id and voice and not text:
                    file_id = voice.get("file_id")
                    if file_id:
                        _send_telegram_message_sync(chat_id, "🎤 Расшифровываю голосовое сообщение...")
                        text = transcribe_voice(file_id)
                        if not text:
                            _send_telegram_message_sync(chat_id, "❌ Не удалось расшифровать голосовое сообщение")
                            continue

                if not chat_id or not text:
                    continue

                # Handle /start deep link binding (D-18)
                if text.startswith("/start ") and chat_id:
                    parts = text.split()
                    if len(parts) > 1:
                        param = parts[1]
                        if param.startswith("bind_"):
                            web_session_id = param[5:]
                            lead_id = _session_bindings.pop(web_session_id, None)
                            if lead_id:
                                _chat_lead_map[chat_id] = lead_id
                                _save_bindings()
                                logger.info(f"Session bound: chat_id={chat_id} -> lead_id={lead_id}")
                                _send_telegram_message_sync(
                                    chat_id,
                                    "✅ Ваш чат привязан к аккаунту AIM. Оператор готов ответить на ваши вопросы."
                                )
                                continue

                # Process through Hermes Operator
                _process_update_sync(message, chat_id, text)

        except Exception as e:
            consecutive_errors += 1
            backoff = min(5 * consecutive_errors, 60)
            logger.error(f"Polling error (attempt {consecutive_errors}, backoff {backoff}s): {e}")
            time.sleep(backoff)

    logger.info("Telegram polling stopped")


def start_polling():
    """Start Telegram polling in a separate thread via run_in_executor.

    Called from the first /health request AFTER uvicorn has bound to port 8000.
    Using run_in_executor (separate OS thread) guarantees Telegram API calls
    never block the uvicorn event loop.
    """
    global _polling_future, _polling_stop
    if _polling_future is not None and not _polling_future.done():
        logger.warning("Polling already running, skipping start")
        return
    _polling_stop = False
    loop = asyncio.get_running_loop()
    _polling_future = loop.run_in_executor(None, _polling_loop_sync)


async def stop_polling():
    """Stop the Telegram polling thread. Called from on_event("shutdown")."""
    global _polling_stop, _polling_future
    _polling_stop = True
    if _polling_future:
        _polling_future.cancel()  # no-op for running threads, but marks cancelled
        try:
            await asyncio.get_running_loop().run_in_executor(
                None, lambda: _polling_future.result(timeout=5) if not _polling_future.done() else None
            )
        except Exception:
            pass
        _polling_future = None
